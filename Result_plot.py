import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches

# Define file paths
file_paths = ["run1.xlsx", "run2.xlsx", "run3.xlsx", "run4.xlsx", "run5.xlsx"]

# Read data from files
data_frames = [pd.read_excel(file, sheet_name="Data", header=None) for file in file_paths]
model_frames = [pd.read_excel(file, sheet_name="Model", header=None) for file in file_paths]

# Function to convert Excel column name to index
def col_name_to_index(col_name):
    col_name = col_name.upper()
    index = 0
    for char in col_name:
        index = index * 26 + (ord(char) - ord('A') + 1)
    return index - 1

# Get and process the data for the title
title1_part1 = data_frames[0].iloc[2, 1].strftime('%Y-%m-%d')
title1_part2 = data_frames[4].iloc[2, 1].strftime('%Y-%m-%d')
fig1_title = f"{title1_part1}--{title1_part2} Susteon Baseline"

# Generate Fig1
fig1_limit = 150
plt.figure(figsize=(10, 6))
max_lines = []
x_offset = 0
for i, df in enumerate(data_frames):
    max_line_no = int(df.iloc[0, 0]) + 1
    max_lines.append(max_line_no)
    if i == 0:
        x_values = df.iloc[4:max_line_no, col_name_to_index('HM')]
        y_values = df.iloc[4:max_line_no, col_name_to_index('BH')]
    else:
        x_values = df.iloc[4:max_line_no, col_name_to_index('HM')] + x_offset
        y_values = df.iloc[4:max_line_no, col_name_to_index('BH')]
    x_offset = x_values.iloc[-1]
    plt.plot(x_values, y_values, label=f"Run {i + 1}")

for line in max_lines[:-1]:
    plt.axvline(x=line, linestyle='--', color='grey')

plt.title(fig1_title)
plt.xlabel("Hours")
plt.ylabel("MEA Conc (wt%)-No Loading")
plt.legend(title="Continuous Titration")
plt.grid(True)
plt.xlim(0, fig1_limit)
fig1_path = 'fig1.png'
plt.savefig(fig1_path)
plt.close()

# Generate Fig2
plt.figure(figsize=(10, 6))
x_values = [df.iloc[max_line_no, col_name_to_index('HW')] for df, max_line_no in zip(data_frames, max_lines)]
y_values = [df.iloc[211, col_name_to_index('E')] for df in model_frames]

plt.scatter(x_values, y_values, color='blue')
plt.plot([1, max(x_values)], [1, max(x_values)], color='black', label='y=x')
plt.plot([1, max(x_values)], [1.05 * x for x in [1, max(x_values)]], linestyle='--', color='blue', label='+5% Deviation')
plt.plot([1, max(x_values)], [0.95 * x for x in [1, max(x_values)]], linestyle='--', color='blue', label='-5% Deviation')
plt.plot([1, max(x_values)], [1.15 * x for x in [1, max(x_values)]], linestyle='--', color='red', label='+15% Deviation')
plt.plot([1, max(x_values)], [0.85 * x for x in [1, max(x_values)]], linestyle='--', color='red', label='-15% Deviation')

plt.title("NGCC Case")
plt.xlabel("PSTU Experimental SRD (GJ/tCO2)")
plt.ylabel("Process Model SRD (GJ/tCO2)")
plt.legend()
plt.grid(True)
fig2_path = 'fig2.png'
plt.savefig(fig2_path)
plt.close()

# Generate Fig3
plt.figure(figsize=(10, 6))
for i, (df, model_df) in enumerate(zip(data_frames, model_frames)):
    max_line_no = int(df.iloc[0, 0]) + 1
    x_values = df.iloc[max_line_no, col_name_to_index('HK')]
    y_values = df.iloc[max_line_no, col_name_to_index('HW')]
    plt.scatter(x_values, y_values, label=f"Run {i + 1}")

model_x_values = [df.iloc[206, col_name_to_index('E')] for df in model_frames]
model_y_values = [df.iloc[211, col_name_to_index('E')] for df in model_frames]
plt.plot(model_x_values, model_y_values, label="Model Curve", color='red')

plt.title("NGCC Case")
plt.xlabel("L/G (mass/mass)")
plt.ylabel("SRD (Gt/tCO2)")
plt.legend()
plt.grid(True)
fig3_path = 'fig3.png'
plt.savefig(fig3_path)
plt.close()

# Create PPTX file
prs = Presentation()

# Add CO2_capture.png image to the first slide
slide_layout = prs.slide_layouts[5]  # Choose a blank layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "CO2 Capture Results"

img_path = 'CO2_capture.png'
slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Add Fig1 to new slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Fig1: MEA Conc (wt%)"

slide.shapes.add_picture(fig1_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Add Fig2 to new slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Fig2: NGCC Case"

slide.shapes.add_picture(fig2_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Add Fig3 to new slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Fig3: NGCC Case"

slide.shapes.add_picture(fig3_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Save PPTX file
prs.save("NCCC campaign results summary.pptx")
