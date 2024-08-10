import pandas as pd
import openpyxl
import numpy as np
import matplotlib.pyplot as plt
from sklearn.linear_model import LinearRegression
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Add a fitted line and return the equation and R² value
def add_trendline(ax, x, y):
    x = np.array(x).reshape(-1, 1)
    y = np.array(y)
    reg = LinearRegression().fit(x, y)
    slope = reg.coef_[0]
    intercept = reg.intercept_
    r2 = reg.score(x, y)
    ax.plot(x, reg.predict(x), color='red', linestyle='--', linewidth=1)
    return f'y = {slope:.2f}x + {intercept:.2f}', r2

# Helper function: convert Excel column name to DataFrame column index
def excel_col_to_num(col):
    num = 0
    for c in col:
        if 'A' <= c <= 'Z':
            num = num * 26 + (ord(c) - ord('A')) + 1
    return num - 1  # Column index starts from 0

# Read the value of last_row_df_extended
def read_last_row(address_path):
    wb = openpyxl.load_workbook(address_path)
    ws = wb['calculation']
    last_row_df_extended = ws['A1'].value
    return last_row_df_extended

# Plot and save charts
def draw_plot(output_path: str, df_final, last_row_df_extended: int):
    # Read the existing workbook
    wb = openpyxl.load_workbook(output_path)

    # Generate charts using Matplotlib and save as images
    charts = [
        ('Liquid Flow', 'HG', 'HM', 'liquid flow rate', 'liquid flow rate (lb/hr)'),
        ('Gas Flow', 'W', 'HM', 'flue gas flow rate', 'gas flow rate (lb/hr)'),
        ('CO2%', ['DF', 'BJ'], 'HM', 'concentration of CO2 in flue gas', 'CO2 (%)', ['outlet', 'inlet']),
        ('FG Temp', ['IE', 'IF'], 'HM', 'flue gas temperature', 'T/C (%)', ['before dilution', 'after dilution']),
        ('washwaterflow', 'DI', 'HM', 'wash water circulation flow', 'wash water flow (lb/hr)'),
        ('washtemp', 'IC', 'HM', 'wash temperature', 'T/C'),
        ('lean temperature', 'ID', 'HM', 'lean solvent temperature', 'T/C'),
        ('Bed location', 'CW', 'HM', 'Number of absorber beds', 'Number of Beds'),
        ('flue gas pressure', 'DP', 'HM', 'flue gas pressure', 'pressure (psig)'),
        ('solvent pressure', 'ET', 'HM', 'inlet solvent pressure', 'pressure (psig)'),
        ('CO2 product pressure', 'FR', 'HM', 'CO2 outlet pressure', 'pressure (psig)'),
        ('CO2temp', 'IB', 'HM', 'CO2 outlet temperature', 'T/C'),
        ('LP steam pressure', 'HD', 'HM', 'LP steam pressure', 'steam pressure (psig)'),
        ('Liq.Fl.Reboiler', 'EZ', 'HM', 'reboiler solvent flow', 'solvent flow (gpm)'),
        ('inter-cooling flow', ['CC', 'CE'], 'HM', 'solvent flow in inter-cooler', 'inter-cooling solvent flow (gpm)', ['HX20401', 'HX20402']),
        ('solvent temperatures', ['IA', 'HZ', 'HY', 'HX'], 'HM', 'solvent temperatures on inter-cooling', 'T/C', ['outlet (HX20402)', 'outlet (HX20401)', 'inlet (HX20402)', 'inlet (HX20401)']),
        ('liquid levels', ['GE', 'GN', 'GO'], 'HM', 'liquid levels of PSTU', 'liquid level (%)', ['adsorber sump', 'condenser sump', 'regenerator sump']),
        ('condensate correction', ['GL', 'HN'], 'HM', 'steam correction', 'condensate (lb/hr)', ['before', 'after']),
        ('SRDs', 'HW', 'HM', 'SRD variation based on condensate', 'SRD (GJ/tCO2)'),
        ('lean approach', ['JU', 'JV'], 'HM', 'lean approach', 'hot approach (T/C)'),
        ('rich approach', ['JT', 'JW'], 'HM', 'rich approach', 'cold side approach (T/C)'),
        ('CO2 loading analysis', ['JY', 'JZ'], 'HM', 'CO2 loading trend', 'CO2 loading (mol CO2/mol MEA)', ['lean', 'rich']),
    ]

    for chart_info in charts:
        chart_name = chart_info[0]
        y_cols = chart_info[1]
        x_col = chart_info[2]
        title = chart_info[3]
        y_axis_title = chart_info[4]
        legends = chart_info[5] if len(chart_info) > 5 else None

        fig, ax = plt.subplots()
        x = df_final.iloc[4:last_row_df_extended + 1, excel_col_to_num(x_col)]

        if isinstance(y_cols, list):
            if legends is not None:
                for y_col, legend in zip(y_cols, legends):
                    y = df_final.iloc[4:last_row_df_extended + 1, excel_col_to_num(y_col)]
                    ax.scatter(x, y, label=legend)
                ax.legend()
            else:
                for y_col in y_cols:
                    y = df_final.iloc[4:last_row_df_extended + 1, excel_col_to_num(y_col)]
                    ax.scatter(x, y)
        else:
            y = df_final.iloc[4:last_row_df_extended + 1, excel_col_to_num(y_cols)]
            ax.scatter(x, y, label=title)
            if chart_name == 'SRDs':
                expression, r2 = add_trendline(ax, x, y)
                ax.annotate(f'Fit Line: {expression}\nR² = {r2:.2f}', xy=(0.05, 0.95), xycoords='axes fraction', fontsize=10,
                            ha='left', va='top', bbox=dict(boxstyle="round,pad=0.3", edgecolor="black", facecolor="white"))
                # Add value of E208
                e212_value = wb['Model']['E212'].value
                ax.annotate(f'SRD: {e212_value}', xy=(0.95, 0.95), xycoords='axes fraction', fontsize=10,
                            ha='right', va='top',
                            bbox=dict(boxstyle="round,pad=0.3", edgecolor="black", facecolor="white"))

        ax.set_title(title)
        ax.set_xlabel('time on stream (hr)')
        ax.set_ylabel(y_axis_title)
        plt.savefig(f'./{chart_name}.png')
        plt.close()

        # Insert image into Excel workbook
        if chart_name in wb.sheetnames:
            wb.remove(wb[chart_name])

        ws = wb.create_sheet(title=chart_name)
        img = openpyxl.drawing.image.Image(f'./{chart_name}.png')
        ws.add_image(img, 'A1')

    # Generate Strip.temp.Profile chart
    fig, ax = plt.subplots()
    data_ws = wb['Data']
    model_ws = wb['Model']
    x_data = [data_ws[f'JJ{last_row_df_extended + 2}'].value, data_ws[f'JK{last_row_df_extended + 2}'].value,
              data_ws[f'JL{last_row_df_extended + 2}'].value, data_ws[f'JM{last_row_df_extended + 2}'].value,
              data_ws[f'JN{last_row_df_extended + 2}'].value, data_ws[f'JO{last_row_df_extended + 2}'].value,
              data_ws[f'JP{last_row_df_extended + 2}'].value, data_ws[f'JQ{last_row_df_extended + 2}'].value,
              data_ws[f'JR{last_row_df_extended + 2}'].value, data_ws[f'JS{last_row_df_extended + 2}'].value]
    y_data = [data_ws[f'JJ{last_row_df_extended + 3}'].value, data_ws[f'JK{last_row_df_extended + 3}'].value,
              data_ws[f'JL{last_row_df_extended + 3}'].value, data_ws[f'JM{last_row_df_extended + 3}'].value,
              data_ws[f'JN{last_row_df_extended + 3}'].value, data_ws[f'JO{last_row_df_extended + 3}'].value,
              data_ws[f'JP{last_row_df_extended + 3}'].value, data_ws[f'JQ{last_row_df_extended + 3}'].value,
              data_ws[f'JR{last_row_df_extended + 3}'].value, data_ws[f'JS{last_row_df_extended + 3}'].value]
    x_model = [model_ws[f'M{i}'].value for i in range(120, 163)]
    y_model = [model_ws[f'L{i}'].value for i in range(120, 163)]

    ax.scatter(x_data, y_data, label='Experiment', color='blue')
    ax.plot(x_model, y_model, label='Simulation', color='red')

    ax.set_title('Stripper Temperature Profiles')
    ax.set_xlabel('T/C')
    ax.set_ylabel('Stripper Height (ft)')
    ax.legend()
    plt.savefig('./Strip.temp.Profile.png')
    plt.close()

    ws = wb.create_sheet(title='Strip.temp.Profile')
    img = openpyxl.drawing.image.Image('./Strip.temp.Profile.png')
    ws.add_image(img, 'A1')

    # Generate Abs.temp.Profile chart
    fig2, ax2 = plt.subplots()
    x_data2 = [data_ws[f'IJ{last_row_df_extended + 2}'].value, data_ws[f'IK{last_row_df_extended + 2}'].value,
               data_ws[f'IL{last_row_df_extended + 2}'].value, data_ws[f'IM{last_row_df_extended + 2}'].value,
               data_ws[f'IN{last_row_df_extended + 2}'].value, data_ws[f'IO{last_row_df_extended + 2}'].value,
               data_ws[f'IP{last_row_df_extended + 2}'].value, data_ws[f'IQ{last_row_df_extended + 2}'].value,
               data_ws[f'IR{last_row_df_extended + 2}'].value, data_ws[f'IS{last_row_df_extended + 2}'].value,
               data_ws[f'IT{last_row_df_extended + 2}'].value, data_ws[f'IU{last_row_df_extended + 2}'].value,
               data_ws[f'IV{last_row_df_extended + 2}'].value, data_ws[f'IW{last_row_df_extended + 2}'].value,
               data_ws[f'IX{last_row_df_extended + 2}'].value, data_ws[f'IY{last_row_df_extended + 2}'].value,
               data_ws[f'IZ{last_row_df_extended + 2}'].value, data_ws[f'JA{last_row_df_extended + 2}'].value,
               data_ws[f'JB{last_row_df_extended + 2}'].value, data_ws[f'JC{last_row_df_extended + 2}'].value,
               data_ws[f'JD{last_row_df_extended + 2}'].value]

    y_data2 = [data_ws[f'IJ{last_row_df_extended + 3}'].value, data_ws[f'IK{last_row_df_extended + 3}'].value,
               data_ws[f'IL{last_row_df_extended + 3}'].value, data_ws[f'IM{last_row_df_extended + 3}'].value,
               data_ws[f'IN{last_row_df_extended + 3}'].value, data_ws[f'IO{last_row_df_extended + 3}'].value,
               data_ws[f'IP{last_row_df_extended + 3}'].value, data_ws[f'IQ{last_row_df_extended + 3}'].value,
               data_ws[f'IR{last_row_df_extended + 3}'].value, data_ws[f'IS{last_row_df_extended + 3}'].value,
               data_ws[f'IT{last_row_df_extended + 3}'].value, data_ws[f'IU{last_row_df_extended + 3}'].value,
               data_ws[f'IV{last_row_df_extended + 3}'].value, data_ws[f'IW{last_row_df_extended + 3}'].value,
               data_ws[f'IX{last_row_df_extended + 3}'].value, data_ws[f'IY{last_row_df_extended + 3}'].value,
               data_ws[f'IZ{last_row_df_extended + 3}'].value, data_ws[f'JA{last_row_df_extended + 3}'].value,
               data_ws[f'JB{last_row_df_extended + 3}'].value, data_ws[f'JC{last_row_df_extended + 3}'].value,
               data_ws[f'JD{last_row_df_extended + 3}'].value]

    # Filter data points based on x_data2
    filtered_x_data2 = []
    filtered_y_data2 = []
    for x, y in zip(x_data2, y_data2):
        if 20 <= x <= 90:  # Only include data points where x is between 20 and 90
            filtered_x_data2.append(x)
            filtered_y_data2.append(y)

    x_model2 = [model_ws[f'M{i}'].value for i in range(10, 102)]
    y_model2 = [model_ws[f'L{i}'].value for i in range(10, 102)]

    ax2.scatter(filtered_x_data2, filtered_y_data2, label='Experiment', color='blue')
    ax2.plot(x_model2, y_model2, label='Simulation', color='red')

    ax2.set_title('Absorber Temperature Profiles')
    ax2.set_xlabel('T/C')
    ax2.set_ylabel('Absorber Column height (ft)')
    ax2.set_xlim(25, 75)  # Set x-axis range from 25 to 75
    ax2.legend()
    plt.savefig('./Abs.temp.Profile.png')
    plt.close()

    ws = wb.create_sheet(title='Abs.temp.Profile')
    img = openpyxl.drawing.image.Image('./Abs.temp.Profile.png')
    ws.add_image(img, 'A1')

    # Generate CO2 capture chart
    fig3, ax3 = plt.subplots()
    x_co2 = [data_ws[f'HM{i}'].value for i in range(5, last_row_df_extended + 3)]
    y_co2_capture = [data_ws[f'GP{i}'].value for i in range(5, last_row_df_extended + 3)]
    y_lg = [data_ws[f'HK{i}'].value for i in range(5, last_row_df_extended + 3)]

    ax3.scatter(x_co2, y_co2_capture, label='CO2 Capture (%)', color='blue')
    ax3.set_xlabel('Time on Stream (hr)')
    ax3.set_ylabel('CO2 Capture (%)', color='blue')
    ax3.tick_params(axis='y', labelcolor='blue')

    ax4 = ax3.twinx()
    ax4.scatter(x_co2, y_lg, label='L/G (mass/mass)', color='red')
    ax4.set_ylabel('L/G (mass/mass)', color='red')
    ax4.tick_params(axis='y', labelcolor='red')

    fig3.suptitle('CO2 Capture')
    fig3.legend(loc='upper right')
    plt.savefig('./CO2_capture.png')
    plt.close()

    ws = wb.create_sheet(title='CO2 capture')
    img = openpyxl.drawing.image.Image('./CO2_capture.png')
    ws.add_image(img, 'A1')

    # Create LbyG worksheet and generate chart
    ws_lbyg = wb.create_sheet(title='LbyG')
    fig4, ax4 = plt.subplots()
    x_lbyg = df_final.iloc[4:last_row_df_extended + 1, excel_col_to_num('HM')]
    y_lbyg = df_final.iloc[4:last_row_df_extended + 1, excel_col_to_num('HK')]

    ax4.scatter(x_lbyg, y_lbyg)
    ax4.set_title('L by G Profiles')
    ax4.set_xlabel('Time on Stream (hr)')
    ax4.set_ylabel('L/G (mass/mass)')
    plt.savefig('./LbyG.png')
    plt.close()

    img_lbyg = openpyxl.drawing.image.Image('./LbyG.png')
    ws_lbyg.add_image(img_lbyg, 'A1')

    wb.save(output_path)
    print("Figures written to Excel successfully.")

def create_pptx_from_excel_images(excel_path: str, pptx_path: str):
    prs = Presentation()

    # Add the first slide
    slide_layout = prs.slide_layouts[0]  # Use title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "NCCC-testing"
    subtitle.text = "Baseline - MEA"

    # Get all PNG files in the current folder
    png_files = [f for f in os.listdir('.') if f.endswith('.png')]

    for image_path in png_files:
        slide_layout = prs.slide_layouts[5]  # Use title and content slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = os.path.splitext(image_path)[0]

        # Insert image into slide
        left = Inches(1)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(image_path, left, top, height=Inches(5))

    # Read data from Excel and add to PPT
    wb = openpyxl.load_workbook(excel_path)
    model_ws = wb['Model']

    slide_layout = prs.slide_layouts[1]  # Use title and content slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Simulation Results"

    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for row in range(184, 220):
        d_value = model_ws[f'D{row}'].value
        e_value = model_ws[f'E{row}'].value
        if d_value is not None and e_value is not None:
            p = text_frame.add_paragraph()
            p.text = f'{d_value} : {e_value}'
            p.font.size = Pt(9)

    prs.save(pptx_path)
    print(f"Presentation saved to {pptx_path}")

# Read last_row_df_extended and df_final
last_row_df_extended = read_last_row('address.xlsx')
df_final = pd.read_pickle('df_final.pkl')

# Plot charts and save to Excel
draw_plot('qiaochu_processed_data.xlsx', df_final, last_row_df_extended)

# Create PPT from images and charts in Excel
create_pptx_from_excel_images('qiaochu_processed_data.xlsx', 'Results-NCCC.pptx')
